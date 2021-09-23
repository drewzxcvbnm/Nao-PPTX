import threading
import win32com.client as win32
from eventmanager import Eventloop, Event, binaryPredicate
from pptx import Presentation
import time
from services import atts, tts, touch, motion
import qi
from slidepresentor import SlidePresentor
from web.webinterface import WebInterface

langs = {
    "rus": "Alyona22Enhanced",
    "eng": "naoenu"
}


class PresentationReader:

    def __init__(self, path, lang="eng"):
        name = path.split('\\')[-1].split('.')[0]
        self.lock = threading.Lock()
        self.presentation_id = WebInterface.create_presentation(name)
        self.stop = False
        tts.setVoice(langs[lang])
        self.path = path
        self.ppoint = win32.gencache.EnsureDispatch('Powerpoint.Application')
        self.ppoint.Visible = True
        self.presentation = self.ppoint.Presentations.Open(path)
        self.slide_show = self.presentation.SlideShowSettings.Run()
        self.event_handler = Eventloop()
        self.event_handler.addEvent(
            Event(self._pause, [], binaryPredicate(lambda: touch.getStatus()[8][1], False, True)))
        self.event_handler.addEvent(
            Event(self._next_slide, [], binaryPredicate(lambda: touch.getStatus()[7][1], False, True)))
        self.event_handler.addEvent(
            Event(self._prev_slide, [], binaryPredicate(lambda: touch.getStatus()[9][1], False, True)))
        self.slide_reader = SlidePresentor(self.slide_show, self.presentation_id)

    def read_slides(self):
        self.i_slide = 0
        self.ppt = Presentation(self.path)
        self.event_handler.start()
        while self.i_slide < len(self.ppt.slides):
            if self.stop:
                break
            self.lock.acquire()
            self.slide_show.View.GotoSlide(self.i_slide + 1)
            self.slide_reader.read_slide(self.ppt.slides[self.i_slide])
            self.i_slide += 1
            self.lock.release()
            time.sleep(1)

    def close(self):
        self.event_handler.stop()
        self.presentation.Close()
        self.ppoint.Quit()
        self.event_handler.join()

    def _pause(self):
        print ("pause")
        tts.stopAll()
        self.lock.acquire()
        while not touch.getStatus()[8][1]:
            time.sleep(0.1)
            print ("111")
        print ("resume")
        self.lock.release()
        time.sleep(1)

    def _stop(self):
        self.stop = True
        tts.stopAll()

    def _next_slide(self):
        tts.stopAll()

    def _prev_slide(self):
        self.i_slide -= 2
        tts.stopAll()
