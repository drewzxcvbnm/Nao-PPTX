import pythoncom
import win32api
import win32com.client as win32
from web.webinterface import WebInterface
from pptx import Presentation as PPTXPresentation
from constants import EVENT_ARG_DELIMITER
from services import mem

# mem event MUST be declared immediately here
mem.declareEvent("event")
subscriber = mem.subscriber("event")


class Presentation:

    def __init__(self, path):
        name = path.split('\\')[-1].split('.')[0]
        self.presentation_id = WebInterface.create_presentation(name)
        self.path = path
        self.com_ppoint = win32.gencache.EnsureDispatch('Powerpoint.Application')
        self.com_presentation = self._open_presentation()
        self.com_slide_show = self.com_presentation.SlideShowSettings.Run()
        self.pptx_presentation = PPTXPresentation(self.path)
        self.slides = self.pptx_presentation.slides
        self.surveys = {}
        self.ongoing_events = []
        self.event_map = {}
        subscriber.signal.connect(self.handle_event)

    def handle_event(self, event):
        eventname = event.split(EVENT_ARG_DELIMITER)[0]
        print("HANDLING " + eventname)
        if eventname not in self.event_map.keys():
            print("ERROR: eventmap cannot handle event:" + eventname)
        mem.insertData("event", None)
        self.event_map[eventname].execute_event(event)

    def __del__(self):
        self.com_slide_show.View.Exit()
        self.com_presentation.Close()
        self.com_ppoint.Quit()
        WebInterface.delete_presentation(self.presentation_id)
        for values in self.event_map.values():
            del values

    def _open_presentation(self):
        try:
            return self.com_ppoint.Presentations.Open(self.path)
        except pythoncom.com_error as err:
            msg = win32api.FormatMessage(err[2][5])
            raise RuntimeError("Error opening presentation: " + msg)
