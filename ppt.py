import win32com.client as win32
import time
import qi

from eventmanager import *
from pptx import Presentation

def main(session):
    global tts
    tts = session.service("ALTextToSpeech")
    atts = session.service("ALAnimatedSpeech")
    alife = session.service("ALAutonomousLife")
    motion = session.service("ALMotion")
    asr = session.service("ALSpeechRecognition")
    touch = session.service("ALTouch")

    # alife.setState("disabled")
    # alife.stopAll()
    asr.pause(1)
    # asr.unsubscribe("WordRecognized")
    # motion.stopMove()
    # motion.wakeUp()


    tts.setVoice("Alyona22Enhanced")
    # tts.setVoice("naoenu")
    # filePath = "C:\\Users\\tsi_nao\\Desktop\\TEXTANNOTATION.pptx"
    filePath = "C:\\Users\\tsi_nao\\Desktop\\TESTTHIS.pptx"
    ppoint = win32.gencache.EnsureDispatch('Powerpoint.Application')
    ppoint.Visible = True
    presentation = ppoint.Presentations.Open(filePath)
    presentation.SlideShowSettings.ShowWithAnimation = False
    slideShow = presentation.SlideShowSettings.Run()


    ppt=Presentation(filePath)
    notes = []

    global eventHandler
    eventHandler = Eventloop()
    eventHandler.addEvent(Event( skipThisShit, [], binaryPredicate(lambda: touch.getStatus()[8][1], True, False)))
    eventHandler.start()


    for page, slide in enumerate(ppt.slides):
        textNote = slide.notes_slide.notes_text_frame.text
        notes.append(textNote)
        print str(notes[page].encode('utf-8'))
        say = qi.async(atts.say,(str(notes[page].encode('utf-8'))),delay=100)
        # tts.say(str(notes[page].encode('utf-8')))
        say.wait()

        slideShow.View.Next()
        time.sleep(1)

    presentation.Close()
    ppoint.Quit()
    # asr.subscribe("WordRecognized")
    eventHandler.stop()
    eventHandler.join()

def skipThisShit():
    print ("help me")
    tts.stopAll()


session = qi.Session()
# session.connect("tcp://" + "192.168.253.155" + ":" + "9559")
session.connect("tcp://" + "192.168.252.226" + ":" + "9559")
main(session)