import win32com.client as win32
from eventmanager import *
from pptx import Presentation
import time
from services import atts, tts, touch
import qi
langs = {
    "rus" : "Alyona22Enhanced",
    "eng" : "naoenu"
}

class PresentationReader:

    def __init__(self, path, lang = "eng"):
        self.stop = False
        tts.setVoice(langs[lang])
        self.path = path
        self.ppoint = win32.gencache.EnsureDispatch('Powerpoint.Application')
        self.ppoint.Visible = True
        self.presentation = self.ppoint.Presentations.Open(path)
        self.presentation.SlideShowSettings.ShowWithAnimation = False
        self.slideShow = self.presentation.SlideShowSettings.Run()
        self.eventHandler = Eventloop()
        self.eventHandler.addEvent(Event(self._toggle, [], binaryPredicate(lambda: touch.getStatus()[8][1], False, True)))
        self.eventHandler.addEvent(Event(self._nextSlide, [], binaryPredicate(lambda: touch.getStatus()[7][1], False, True)))
        self.eventHandler.addEvent(Event(self._prevSlide, [], binaryPredicate(lambda: touch.getStatus()[9][1], False, True)))

    def readSlides(self):
        self.ppt=Presentation(self.path)
        self.iSlide = 0
        self.eventHandler.start()
        while self.iSlide < len(self.ppt.slides):
            while self.stop:
                time.sleep(0.5)
            self._readSlide(self.iSlide)
            self.iSlide += 1
            self.slideShow.View.GotoSlide(self.iSlide + 1)
            time.sleep(1)

    def close(self):
        self.eventHandler.stop()
        self.presentation.Close()
        self.ppoint.Quit()
        self.eventHandler.join()

    def _readSlide(self, i):
        global atts
        slide = self.ppt.slides[i]
        textNote = slide.notes_slide.notes_text_frame.text
        notes = textNote.encode('utf-8')
        print str(notes)
        say = qi.async(atts.say, (str(notes)), delay=100)
        # tts.say(str(notes[page].encode('utf-8')))
        say.wait()

    def _toggle(self):
        self.stop = not self.stop
        tts.stopAll()
   
    def _nextSlide(self):
        global tts
        tts.stopAll()

    def _prevSlide(self):
        global tts
        self.iSlide -= 2
        tts.stopAll()
